"""
Generates small, deterministic sample source documents for every input
type AEGIS is required to ingest (Requirement.pdf/Q&A.pdf gap-fill plan,
Phase 5). Run once:

    python tests/sample_data/generate_samples.py

Files are built programmatically (no flaky/slow external downloads needed
on an air-gapped or bandwidth-limited machine) but are representative of
real-world content: an ICS component datasheet, a security advisory, a
corporate asset database, firmware source code with communication
interfaces and a possible hardcoded secret, a wiring schematic image with
OCR-able labels and connecting lines, and a scanned "handwritten note"
image. Regenerating is idempotent - existing files are overwritten.
"""

from __future__ import annotations

import sqlite3
import zipfile
from pathlib import Path

HERE = Path(__file__).parent


# ── PDF (datasheet) ────────────────────────────────────────────────────────
def make_pdf():
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas

    path = HERE / "datasheet_stm32f407.pdf"
    c = canvas.Canvas(str(path), pagesize=A4)
    c.setFont("Helvetica-Bold", 16)
    c.drawString(50, 780, "STM32F407VGT6 - Datasheet Extract")
    c.setFont("Helvetica", 11)
    lines = [
        "Manufacturer: STMicroelectronics",
        "Part number: STM32F407VGT6",
        "Core: ARM Cortex-M4, 168 MHz, hardware FPU",
        "Communication interfaces: UART, SPI, I2C, CAN bus, USB, Ethernet",
        "Firmware version supported: v2.3.1",
        "Known issue: CVE-2023-44487 (rapid reset / resource exhaustion class)",
        "This device is commonly deployed in industrial controllers using",
        "the Modbus and CANopen protocols for field-bus communication.",
        "",
        "Page 2 continues with electrical characteristics and pin mapping.",
    ]
    y = 750
    for line in lines:
        c.drawString(50, y, line)
        y -= 20
    c.showPage()
    c.setFont("Helvetica-Bold", 14)
    c.drawString(50, 780, "Page 2 - Pinout & Electrical Characteristics")
    c.setFont("Helvetica", 11)
    c.drawString(50, 750, "Operating voltage: 1.8V - 3.6V. Max clock: 168MHz.")
    c.drawString(50, 730, "See schematic diagram (sample_schematic.png) for wiring.")
    c.showPage()
    c.save()
    return path


# ── DOCX (manual) ──────────────────────────────────────────────────────────
def make_docx():
    from docx import Document

    path = HERE / "manual_plc_operations.docx"
    doc = Document()
    doc.add_heading("Siemens S7-1500 PLC - Operations Manual (extract)", level=1)
    doc.add_paragraph(
        "This manual covers safe operating procedures for the Siemens S7-1500 "
        "programmable logic controller, part number 6ES7515-2AM01-0AB0."
    )
    doc.add_paragraph(
        "The controller communicates via PROFINET and supports OPC-UA for "
        "supervisory control integration. Firmware version 2.9.3 is required "
        "for the latest security patches, addressing CVE-2022-38465."
    )
    doc.add_heading("Manufacturer", level=2)
    doc.add_paragraph("Siemens AG, Germany.")
    doc.save(str(path))
    return path


# ── PPTX (briefing slides) ──────────────────────────────────────────────────
def make_pptx():
    from pptx import Presentation
    from pptx.util import Inches

    path = HERE / "briefing_ics_overview.pptx"
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "ICS Component Overview"
    slide1.placeholders[1].text = "Rockwell Allen-Bradley PLCs - assessment briefing"

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Key findings"
    body = slide2.placeholders[1].text_frame
    body.text = "Manufacturer: Rockwell Automation (Allen-Bradley)"
    body.add_paragraph().text = "Protocol: EtherNet/IP, Modbus TCP"
    body.add_paragraph().text = "Advisory: CVE-2021-22681 affects authentication"
    prs.save(str(path))
    return path


# ── XLSX (asset register) ───────────────────────────────────────────────────
def make_xlsx():
    import openpyxl

    path = HERE / "asset_register.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Assets"
    ws.append(["Asset ID", "Manufacturer", "Part Number", "Protocol", "Firmware Version"])
    ws.append(["A-001", "Schneider Electric", "TM251MESE", "Modbus", "v5.1.2"])
    ws.append(["A-002", "Honeywell", "ControlEdge-900", "BACnet", "v3.0.0"])
    ws.append(["A-003", "Omron", "NX1P2-1140DT1", "EtherCAT", "v1.4.0"])
    wb.save(str(path))
    return path


# ── CSV / TXT / MD ───────────────────────────────────────────────────────────
def make_text_files():
    (HERE / "forum_thread.txt").write_text(
        "Forum post: has anyone seen the ESP32-WROOM-32 module drop off "
        "Wi-Fi under load? We're running firmware v1.0.4 and using MQTT "
        "to a broker on the LAN. Suspect a possible attack surface if the "
        "broker is exposed - we found api_key = \"letmein1234\" hardcoded "
        "in one build, which needs fixing.\n",
        encoding="utf-8",
    )
    (HERE / "release_notes.md").write_text(
        "# Release Notes\n\n"
        "## v2.1.0\n"
        "- Added Bluetooth Low Energy (BLE) support for the Nordic Semiconductor "
        "nRF52840 module.\n"
        "- Fixed CVE-2023-1234 buffer overflow in the OTA updater.\n\n"
        "## v2.0.0\n"
        "- Initial release with UART and SPI driver support.\n",
        encoding="utf-8",
    )
    (HERE / "sensor_log.csv").write_text(
        "timestamp,sensor_id,value,unit\n"
        "2026-01-01T00:00:00,TEMP-01,21.4,C\n"
        "2026-01-01T00:05:00,TEMP-01,21.6,C\n"
        "2026-01-01T00:10:00,PRESS-02,101.3,kPa\n",
        encoding="utf-8",
    )
    (HERE / "advisory.html").write_text(
        "<html><head><title>Security Advisory</title></head><body>"
        "<h1>Security Advisory ADV-2026-004</h1>"
        "<p>Affected product: Beckhoff CX5140 embedded PC.</p>"
        "<p>Vulnerability: CVE-2024-9999 allows unauthenticated access via "
        "an exposed EtherCAT diagnostic port.</p>"
        "<p>Manufacturer: Beckhoff Automation, Germany.</p>"
        "</body></html>",
        encoding="utf-8",
    )


# ── Corporate database (SQLite) ─────────────────────────────────────────────
def make_sqlite():
    path = HERE / "corporate_assets.db"
    if path.exists():
        path.unlink()
    conn = sqlite3.connect(str(path))
    cur = conn.cursor()
    cur.execute(
        "CREATE TABLE devices (id INTEGER PRIMARY KEY, manufacturer TEXT, "
        "part_number TEXT, protocol TEXT, install_site TEXT)"
    )
    cur.executemany(
        "INSERT INTO devices (manufacturer, part_number, protocol, install_site) VALUES (?,?,?,?)",
        [
            ("Mitsubishi Electric", "FX5U-32MT", "CC-Link", "Plant A - Line 3"),
            ("Yokogawa", "CENTUM-VP", "HART", "Plant B - Control Room"),
            ("ABB", "AC500-eCo", "Profibus", "Plant A - Line 1"),
        ],
    )
    cur.execute(
        "CREATE TABLE advisories (id INTEGER PRIMARY KEY, cve TEXT, severity TEXT, device_id INTEGER)"
    )
    cur.executemany(
        "INSERT INTO advisories (cve, severity, device_id) VALUES (?,?,?)",
        [("CVE-2020-10001", "High", 1), ("CVE-2019-6812", "Medium", 2)],
    )
    conn.commit()
    conn.close()
    return path


def make_sql_dump():
    path = HERE / "legacy_export.sql"
    path.write_text(
        "CREATE TABLE users (id INT, username VARCHAR(50), role VARCHAR(20));\n"
        "INSERT INTO users VALUES (1, 'admin', 'administrator');\n"
        "INSERT INTO users VALUES (2, 'operator1', 'operator');\n"
        "CREATE TABLE devices (id INT, manufacturer VARCHAR(50), part_number VARCHAR(50));\n"
        "INSERT INTO devices VALUES (1, 'NXP', 'LPC1768FBD100');\n"
        "INSERT INTO devices VALUES (2, 'Renesas', 'RX65N');\n",
        encoding="utf-8",
    )
    return path


# ── Code archive (firmware source) ──────────────────────────────────────────
def make_code_archive():
    src_dir = HERE / "_code_src"
    src_dir.mkdir(exist_ok=True)
    (src_dir / "main.c").write_text(
        "// Firmware entry point\n"
        "#include \"uart_driver.h\"\n"
        "#include \"i2c_driver.h\"\n\n"
        "int main(void) {\n"
        "    uart_init(115200);      // UART comms\n"
        "    i2c_init();             // I2C sensor bus\n"
        "    bootloader_check_ota(); // OTA update path\n"
        "    return 0;\n"
        "}\n",
        encoding="utf-8",
    )
    (src_dir / "server.py").write_text(
        "from flask import Flask\n"
        "app = Flask(__name__)\n\n"
        "API_KEY = \"sk_live_hardcoded_12345\"  # TODO: move to secrets store\n\n"
        "@app.route('/status')\n"
        "def status():\n"
        "    return {'ok': True}\n\n"
        "if __name__ == '__main__':\n"
        "    app.run(host='0.0.0.0', port=8080)  # open network listener\n",
        encoding="utf-8",
    )
    (src_dir / "requirements.txt").write_text("flask==3.0.0\nrequests==2.31.0\n", encoding="utf-8")

    zip_path = HERE / "firmware_src.zip"
    with zipfile.ZipFile(zip_path, "w") as zf:
        for f in src_dir.iterdir():
            zf.write(f, arcname=f.name)
    return zip_path, src_dir


# ── Images: schematic + handwritten annotation + format variants ───────────
def make_images():
    from PIL import Image, ImageDraw, ImageFont

    def _font(size):
        try:
            return ImageFont.truetype("arial.ttf", size)
        except Exception:
            return ImageFont.load_default()

    # Schematic-style image: boxes with reference labels + connecting lines,
    # so the OCR + Hough-line heuristic (services/schematic_vision.py) has
    # something realistic to detect.
    img = Image.new("RGB", (640, 400), "white")
    draw = ImageDraw.Draw(img)
    boxes = [("U1 MCU", (60, 60, 180, 130)), ("U2 SENSOR", (420, 60, 560, 130)),
             ("J1 CONN", (60, 260, 180, 330)), ("U3 PSU", (420, 260, 560, 330))]
    for label, (x1, y1, x2, y2) in boxes:
        draw.rectangle([x1, y1, x2, y2], outline="black", width=2)
        draw.text((x1 + 10, y1 + 30), label, fill="black", font=_font(16))
    # Connecting lines between box centres (candidate wires)
    draw.line([(180, 95), (420, 95)], fill="black", width=3)
    draw.line([(120, 130), (120, 260)], fill="black", width=3)
    draw.line([(500, 130), (500, 260)], fill="black", width=3)
    img.save(HERE / "sample_schematic.png")

    # Scanned "handwritten annotation" style image: printed text plus a
    # squiggly overlay mark that a plain OCR pass will read at very low
    # confidence - exercising services/ingest.py's handwriting heuristic.
    img2 = Image.new("RGB", (640, 300), "white")
    draw2 = ImageDraw.Draw(img2)
    draw2.text((30, 30), "Inspection sheet - Unit 7, Panel B", fill="black", font=_font(18))
    draw2.text((30, 70), "Checked by QA on 2026-06-01. Status: PASS.", fill="black", font=_font(14))
    # A wobbly freehand-style scribble to emulate a handwritten margin note.
    import math
    pts = [(300 + i, 180 + int(14 * math.sin(i / 6))) for i in range(0, 220, 4)]
    draw2.line(pts, fill=(40, 40, 160), width=3)
    draw2.text((300, 210), "chk!", fill=(40, 40, 160), font=_font(16))
    img2.save(HERE / "sample_scanned_note.png")

    # Format variants for the same content, to exercise every accepted
    # image suffix end-to-end.
    img.convert("RGB").save(HERE / "sample_photo.jpg", quality=85)
    img.convert("RGB").save(HERE / "sample_photo.bmp")
    img.convert("RGB").save(HERE / "sample_photo.tiff")


def main():
    print("Generating AEGIS sample test data in", HERE)
    make_pdf()
    make_docx()
    make_pptx()
    make_xlsx()
    make_text_files()
    make_sqlite()
    make_sql_dump()
    make_code_archive()
    make_images()
    print("Done. Files:")
    for f in sorted(HERE.iterdir()):
        if f.is_file():
            print(" -", f.name, f"({f.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main()
