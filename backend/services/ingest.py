"""
Ingest pipeline
  1. Detect file type and extract raw text
  2. Split into overlapping chunks
  3. Generate embeddings and store in ChromaDB
  4. Index in Whoosh full-text engine
  5. Run NER and store entities
"""

import hashlib
import json
import re
import time
import uuid
from pathlib import Path
from typing import List, Tuple

import chromadb
from whoosh import index as whoosh_index
from whoosh.fields import Schema, TEXT, ID, NUMERIC, KEYWORD
from whoosh.qparser import QueryParser, MultifieldParser
from whoosh.writing import AsyncWriter

from config import (
    CHROMA_DIR, WHOOSH_DIR, UPLOADS_DIR,
    MAX_CHUNK_TOKENS, CHUNK_OVERLAP,
)
from database import get_db

# ── Whoosh schema ─────────────────────────────────────────────────────────────

WHOOSH_SCHEMA = Schema(
    chunk_id      = ID(stored=True, unique=True),
    doc_id        = ID(stored=True),
    collection_id = ID(stored=True),
    filename      = TEXT(stored=True),
    title         = TEXT(stored=True, field_boost=3.0),
    body          = TEXT(stored=False),
    manufacturer  = TEXT(stored=True, field_boost=2.0),
    part_number   = TEXT(stored=True, field_boost=2.5),
    doc_type      = KEYWORD(stored=True),
    page_number   = NUMERIC(stored=True),
)


def _get_whoosh_index():
    p = str(WHOOSH_DIR)
    if whoosh_index.exists_in(p):
        return whoosh_index.open_dir(p)
    return whoosh_index.create_in(p, WHOOSH_SCHEMA)


def _get_chroma():
    settings = chromadb.config.Settings(
        anonymized_telemetry=False,
        chroma_product_telemetry_impl="services.telemetry.NoopTelemetry",
    )
    return chromadb.PersistentClient(path=str(CHROMA_DIR), settings=settings)


# ── Image OCR helpers ─────────────────────────────────────────────────────────

def _find_tesseract() -> str | None:
    """Return path to tesseract binary, or None if not found."""
    import os
    import shutil
    env_path = os.environ.get("AEGIS_TESSERACT_CMD")
    if env_path and Path(env_path).exists():
        return env_path
    if shutil.which("tesseract"):
        return shutil.which("tesseract")
    candidates = [
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
        r"C:\Tesseract-OCR\tesseract.exe",
        "/usr/bin/tesseract",
        "/usr/local/bin/tesseract",
    ]
    for c in candidates:
        if Path(c).exists():
            return c
    return None


def _handwriting_candidates_from_image(img, tess_path: str, known_text: str) -> str:
    """Best-effort handwritten-annotation detection for a single image.

    Handwriting recognition from a flat image is genuinely hard (per
    HMGCC Q&A Q22/Q65). Rather than claiming full handwriting OCR, this
    re-runs Tesseract in sparse-text mode (better suited to scattered
    annotations than a full paragraph read) and flags low-confidence words
    not already present in the primary OCR pass - low OCR confidence is a
    reasonable, explainable proxy for "not printed text", while still being
    honest that this is a heuristic, not verified handwriting recognition.
    """
    import pytesseract

    try:
        data = pytesseract.image_to_data(img, config="--psm 11", output_type=pytesseract.Output.DICT)
    except Exception:
        return ""

    known_tokens = set(re.findall(r"[a-z0-9]+", known_text.lower()))
    candidates = []
    n = len(data.get("text", []))
    for i in range(n):
        word = (data["text"][i] or "").strip()
        try:
            conf = float(data.get("conf", ["-1"])[i])
        except (ValueError, TypeError):
            conf = -1.0
        if word and 0 <= conf < 55 and word.lower() not in known_tokens and len(word) > 1:
            candidates.append(word)

    if not candidates:
        return ""
    return (
        "\n[Possible handwritten annotation candidates - low OCR confidence, "
        f"sparse layout, best-effort only]: {' '.join(candidates[:60])}\n"
    )


def _extract_image_text(file_path: Path) -> str:
    """
    Extract text from an image file.
    Uses Tesseract OCR when available; falls back to image metadata.
    """
    from PIL import Image
    import pytesseract

    img = Image.open(str(file_path))
    width, height = img.size
    mode = img.mode
    fmt = img.format or file_path.suffix.upper().lstrip(".")

    tess_path = _find_tesseract()
    if tess_path:
        pytesseract.pytesseract.tesseract_cmd = tess_path
        # Convert to RGB for best OCR compatibility
        if mode not in ("RGB", "L"):
            img = img.convert("RGB")
        ocr_text = pytesseract.image_to_string(img, config="--psm 3").strip()
        if ocr_text:
            header = (
                f"[Image: {file_path.name} | {fmt} | {width}x{height}px | "
                f"OCR extracted text below]\n\n"
            )
            handwriting = _handwriting_candidates_from_image(img, tess_path, ocr_text)
            return header + ocr_text + handwriting

    # Tesseract not available or OCR produced nothing - return rich metadata
    return (
        f"[Image file: {file_path.name}]\n"
        f"Format: {fmt}\n"
        f"Dimensions: {width} x {height} pixels\n"
        f"Colour mode: {mode}\n"
        f"Note: Install Tesseract OCR (https://github.com/UB-Mannheim/tesseract/wiki) "
        f"to enable full text extraction from images."
    )


# ── Text extraction ───────────────────────────────────────────────────────────

def _extract_pdf_handwriting_candidates(pdf, max_pages: int = 20) -> str:
    """Diff each page's embedded text layer against a re-OCR of the
    rendered page image to isolate probable handwritten overlay annotations
    (HMGCC Q&A Q22: overlaid annotations on an existing document are the
    priority case, ahead of standalone handwritten notes). Best-effort:
    silently degrades to nothing if the PDF can't be rasterised."""
    tess_path = _find_tesseract()
    if not tess_path:
        return ""

    import pytesseract
    pytesseract.pytesseract.tesseract_cmd = tess_path

    out_lines = []
    for page in pdf.pages[:max_pages]:
        try:
            embedded = (page.extract_text() or "")
            embedded_tokens = set(re.findall(r"[a-z0-9]+", embedded.lower()))
            pil_img = page.to_image(resolution=150).original
            ocr_text = pytesseract.image_to_string(pil_img, config="--psm 11")
        except Exception:
            continue

        extra_lines = []
        for line in ocr_text.splitlines():
            line = line.strip()
            if not line:
                continue
            line_tokens = set(re.findall(r"[a-z0-9]+", line.lower()))
            if line_tokens and not (line_tokens & embedded_tokens):
                extra_lines.append(line)
        if extra_lines:
            out_lines.append(
                f"[Possible handwritten annotation candidates on page {page.page_number} - "
                f"best-effort, text present in a page image scan but not the PDF's embedded "
                f"text layer]: {' | '.join(extra_lines[:20])}"
            )
    return ("\n\n" + "\n".join(out_lines)) if out_lines else ""


def _extract_code_archive(file_path: Path) -> Tuple[str, int]:
    """Extract and analyse a small code archive (.zip): reads each member's
    text and runs the code/architecture-insight heuristics (Q14) on it."""
    import zipfile
    from services.code_analysis import is_code_file, analyse_source, format_report

    parts = [f"[Code archive: {file_path.name}]"]
    n_files = 0
    try:
        with zipfile.ZipFile(str(file_path)) as zf:
            for info in zf.infolist():
                if info.is_dir() or info.file_size > 2_000_000:
                    continue
                name = Path(info.filename).name
                suffix = Path(info.filename).suffix
                try:
                    raw = zf.read(info.filename).decode("utf-8", errors="replace")
                except Exception:
                    continue
                n_files += 1
                parts.append(f"\n--- {info.filename} ---\n{raw}")
                if is_code_file(suffix, name):
                    report = analyse_source(raw, name)
                    parts.append(format_report(report, name))
    except Exception as exc:
        return f"[Archive extraction error: {exc}]", 1

    return "\n".join(parts), max(1, n_files)


def _extract_text(file_path: Path, mime: str) -> Tuple[str, int]:
    """Return (text, page_count). Always returns something."""
    suffix = file_path.suffix.lower()
    text = ""
    pages = 1

    try:
        if suffix == ".pdf":
            import pdfplumber
            with pdfplumber.open(str(file_path)) as pdf:
                pages = len(pdf.pages)
                parts = []
                for p in pdf.pages:
                    t = p.extract_text() or ""
                    parts.append(t)
                text = "\n".join(parts)
                text += _extract_pdf_handwriting_candidates(pdf)

        elif suffix in (".docx", ".doc"):
            from docx import Document
            doc = Document(str(file_path))
            text = "\n".join(p.text for p in doc.paragraphs)

        elif suffix in (".pptx", ".ppt"):
            from pptx import Presentation
            prs = Presentation(str(file_path))
            pages = max(1, len(prs.slides))
            parts = []
            for i, slide in enumerate(prs.slides, start=1):
                parts.append(f"[Slide {i}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
            text = "\n".join(parts)

        elif suffix in (".html", ".htm"):
            from bs4 import BeautifulSoup
            import trafilatura
            raw = file_path.read_text(errors="replace")
            extracted = trafilatura.extract(raw)
            text = extracted if extracted else BeautifulSoup(raw, "lxml").get_text(separator="\n")

        elif suffix in (".txt", ".md", ".csv"):
            text = file_path.read_text(errors="replace")

        elif suffix in (".xlsx", ".xls"):
            import openpyxl
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    line = "\t".join(str(c) if c is not None else "" for c in row)
                    parts.append(line)
            text = "\n".join(parts)

        elif suffix in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".webp"):
            text = _extract_image_text(file_path)

        elif suffix in (".db", ".sqlite", ".sqlite3", ".sql"):
            from services.db_ingest import extract_database_text
            text, pages = extract_database_text(file_path)

        elif suffix == ".zip":
            text, pages = _extract_code_archive(file_path)

        else:
            from services.code_analysis import is_code_file, analyse_source, format_report
            if is_code_file(suffix, file_path.name):
                raw = file_path.read_text(errors="replace")
                report = analyse_source(raw, file_path.name)
                text = raw + "\n\n" + format_report(report, file_path.name)
            else:
                text = file_path.read_text(errors="replace")

    except Exception as e:
        text = f"[Extraction error: {e}]"

    return text.strip(), pages


# ── Chunking ──────────────────────────────────────────────────────────────────

def _chunk_text(text: str, max_tokens: int = MAX_CHUNK_TOKENS,
                overlap: int = CHUNK_OVERLAP) -> List[str]:
    """Sentence-aware chunking with overlap.

    Falls back to a fixed-size sliding window for any "sentence" that on its
    own exceeds max_tokens (e.g. unpunctuated text, minified code, OCR output
    with no full stops) so no chunk ever grows unbounded.
    """
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks: List[str] = []
    current: List[str] = []
    current_len = 0

    def _flush():
        nonlocal current, current_len
        if current:
            chunks.append(" ".join(current))
            overlap_words = current[-overlap:] if len(current) > overlap else current
            current = overlap_words[:]
            current_len = len(current)

    for sent in sentences:
        words = sent.split()
        if not words:
            continue

        if len(words) > max_tokens:
            # Oversized "sentence" - flush what we have, then slide a fixed
            # window across it directly so every chunk stays within bounds.
            _flush()
            step = max_tokens - overlap if max_tokens > overlap else max_tokens
            for i in range(0, len(words), step):
                chunks.append(" ".join(words[i:i + max_tokens]))
            current = []
            current_len = 0
            continue

        if current_len + len(words) > max_tokens and current:
            _flush()
        current.extend(words)
        current_len += len(words)

    if current:
        chunks.append(" ".join(current))

    return [c for c in chunks if len(c.strip()) > 20]


# ── NER ───────────────────────────────────────────────────────────────────────

_CVE_RE  = re.compile(r'CVE-\d{4}-\d{4,7}', re.IGNORECASE)
# Broad alphanumeric token scan - a "part number" is any token that mixes
# letters and digits (optionally with hyphens), e.g. STM32F407VGT6,
# ATmega328P, 74HC595, ESP32-WROOM-32. CVE ids are excluded explicitly
# since they also mix letters/digits/hyphens.
_PART_TOKEN_RE = re.compile(r'\b[A-Z0-9][A-Z0-9\-]{3,}[A-Z0-9]\b')
_VER_RE  = re.compile(r'\bv?\d+\.\d+(?:\.\d+)*\b')
_HANDWRITTEN_RE = re.compile(r'\[Possible handwritten annotation[^\]]*\]:\s*(.+)')


def _looks_like_part_number(token: str) -> bool:
    if _CVE_RE.match(token):
        return False
    has_letter = any(c.isalpha() for c in token)
    has_digit = any(c.isdigit() for c in token)
    return has_letter and has_digit and len(token) >= 5

ICS_MANUFACTURERS = {
    "siemens", "schneider", "allen-bradley", "rockwell", "honeywell",
    "abb", "emerson", "yokogawa", "mitsubishi", "omron", "beckhoff",
    "texas instruments", "nxp", "renesas", "stmicroelectronics",
    "microchip", "infineon", "broadcom", "qualcomm", "espressif",
    "nordic semiconductor", "cypress",
}

ICS_PROTOCOLS = {
    "modbus", "profinet", "profibus", "canopen", "can bus", "dnp3",
    "iec 61850", "bacnet", "hart", "foundation fieldbus", "ethercat",
    "powerlink", "sercos", "cc-link", "devicenet", "controlnet",
    "mqtt", "opc-ua", "opc ua",
}


def _extract_entities(text: str) -> List[dict]:
    entities = []
    seen: set = set()

    def _add(etype: str, value: str, context: str = "", confidence: float = 0.8):
        key = (etype, value.upper())
        if key not in seen:
            seen.add(key)
            entities.append({
                "id": str(uuid.uuid4()),
                "entity_type": etype,
                "value": value,
                "confidence": confidence,
                "context": context[:200],
            })

    # CVEs - near-certain given the strict pattern
    for m in _CVE_RE.finditer(text):
        start = max(0, m.start() - 60)
        _add("CVE", m.group(), text[start: m.end() + 60], confidence=0.98)

    # Part numbers (heuristic alphanumeric token scan)
    for m in _PART_TOKEN_RE.finditer(text):
        val = m.group()
        if _looks_like_part_number(val):
            start = max(0, m.start() - 40)
            _add("PART_NUMBER", val, text[start: m.end() + 40], confidence=0.75)

    # Firmware versions - heuristic, requires nearby keyword
    for m in _VER_RE.finditer(text):
        start = max(0, m.start() - 40)
        ctx = text[start: m.end() + 40].lower()
        if any(k in ctx for k in ("firmware", "fw", "version", "ver", "release")):
            _add("FIRMWARE_VERSION", m.group(), text[start: m.end() + 40], confidence=0.65)

    # Manufacturers (case-insensitive keyword scan against a known list)
    lower = text.lower()
    for mfr in ICS_MANUFACTURERS:
        if mfr in lower:
            idx = lower.find(mfr)
            _add("MANUFACTURER", mfr.title(), text[max(0, idx - 40): idx + len(mfr) + 40],
                 confidence=0.85)

    # Protocols (keyword scan against a known list)
    for proto in ICS_PROTOCOLS:
        if proto in lower:
            idx = lower.find(proto)
            _add("PROTOCOL", proto.upper(), text[max(0, idx - 40): idx + len(proto) + 40],
                 confidence=0.85)

    # Handwritten annotation candidates flagged during text extraction
    for m in _HANDWRITTEN_RE.finditer(text):
        content = m.group(1).strip()
        if content:
            _add("HANDWRITTEN_ANNOTATION", content[:150], text[max(0, m.start() - 20): m.end() + 20],
                 confidence=0.45)

    return entities


# ── Embeddings via Ollama ─────────────────────────────────────────────────────

def remove_document_from_indexes(doc_id: str, collection_id: str | None = None) -> None:
    """Remove a document's vectors (Chroma) and keyword entries (Whoosh).

    Called whenever a document row is deleted so the vector/keyword indexes
    never accumulate orphaned entries for documents that no longer exist.
    """
    # Chroma: the collection name is derived from collection_id, but we also
    # sweep the default collection in case it was indexed before a
    # collection-specific one existed.
    try:
        chroma = _get_chroma()
        col_names = set()
        if collection_id:
            col_names.add(f"col_{collection_id.replace('-', '')[:40]}")
        col_names.add("aegis_default")
        for name in col_names:
            try:
                col = chroma.get_collection(name)
                col.delete(where={"doc_id": {"$eq": doc_id}})
            except Exception:
                continue
    except Exception:
        pass

    # Whoosh: delete all chunk documents for this doc_id.
    try:
        ix = _get_whoosh_index()
        writer = AsyncWriter(ix)
        writer.delete_by_term("doc_id", doc_id)
        writer.commit()
    except Exception:
        pass


def _embed_one(text: str) -> List[float]:
    """Embed a single text via Ollama. Returns zero vector on failure."""
    import httpx
    from config import OLLAMA_BASE_URL, EMBEDDING_MODEL
    try:
        r = httpx.post(
            f"{OLLAMA_BASE_URL}/api/embeddings",
            json={"model": EMBEDDING_MODEL, "prompt": text[:4096]},
            timeout=60.0,
        )
        r.raise_for_status()
        return r.json()["embedding"]
    except Exception:
        return [0.0] * 768


def _embed_texts(texts: List[str], max_workers: int = 8) -> List[List[float]]:
    """Embed texts via Ollama using a thread pool for concurrency."""
    from concurrent.futures import ThreadPoolExecutor, as_completed

    results: List[List[float]] = [None] * len(texts)
    with ThreadPoolExecutor(max_workers=max_workers) as pool:
        future_to_idx = {pool.submit(_embed_one, t): i for i, t in enumerate(texts)}
        for future in as_completed(future_to_idx):
            results[future_to_idx[future]] = future.result()
    return results


# ── Main ingest function ──────────────────────────────────────────────────────

def ingest_document(doc_id: str) -> None:
    """Full pipeline: extract → chunk → embed → index → NER. Called in thread."""
    db = get_db()

    def _set_status(status: str, error: str = None):
        db.execute(
            "UPDATE documents SET ingest_status=?, ingest_error=? WHERE id=?",
            (status, error, doc_id),
        )
        db.commit()

    try:
        row = db.execute("SELECT * FROM documents WHERE id=?", (doc_id,)).fetchone()
        if not row:
            return

        _set_status("processing")
        fp = Path(row["file_path"])

        # 1. Extract text
        text, pages = _extract_text(fp, row["mime_type"] or "")
        word_count = len(text.split())

        db.execute(
            "UPDATE documents SET page_count=?, word_count=? WHERE id=?",
            (pages, word_count, doc_id),
        )

        # 1b. Language detection for cultural/bias coverage tracking
        # (desirable requirement - see services/bias.py).
        try:
            from services.bias import tag_document_language
            tag_document_language(doc_id, text)
        except Exception:
            pass

        # 1c. Best-effort spatial grounding for schematic-type images
        # (desirable requirement - see services/schematic_vision.py).
        suffix = fp.suffix.lower()
        if row["doc_type"] == "schematic" and suffix in (
            ".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp"
        ):
            try:
                tess_path = _find_tesseract()
                if tess_path:
                    from services.schematic_vision import analyse_schematic
                    graph = analyse_schematic(fp, tess_path)
                    if graph.get("labels"):
                        db.execute(
                            "INSERT INTO annotations VALUES (?,?,?,?,?,?,?)",
                            (str(uuid.uuid4()), doc_id, None, "schematic_graph",
                             json.dumps(graph), "blue", int(time.time())),
                        )
            except Exception:
                pass

        # 1d. Architecture insight extraction for code files/manifests
        # (essential requirement - see services/code_analysis.py).
        try:
            from services.code_analysis import is_code_file, analyse_source
            if is_code_file(suffix, fp.name):
                raw = fp.read_text(errors="replace")
                report = analyse_source(raw, fp.name)
                if any(report.values()):
                    db.execute(
                        "INSERT INTO annotations VALUES (?,?,?,?,?,?,?)",
                        (str(uuid.uuid4()), doc_id, None, "architecture_insight",
                         json.dumps(report), "green", int(time.time())),
                    )
        except Exception:
            pass

        # 2. Chunk
        chunks = _chunk_text(text)

        # 3. Embeddings
        embeddings = _embed_texts(chunks)

        # 4. Store chunks + embeddings in ChromaDB
        chroma = _get_chroma()
        col_name = f"col_{row['collection_id'].replace('-', '')[:40]}"
        try:
            chroma_col = chroma.get_or_create_collection(
                name=col_name,
                metadata={"hnsw:space": "cosine"},
            )
        except Exception:
            chroma_col = chroma.get_or_create_collection(name="aegis_default")

        chunk_rows = []
        chroma_ids, chroma_docs, chroma_metas, chroma_embeds = [], [], [], []

        for i, (chunk_text, emb) in enumerate(zip(chunks, embeddings)):
            cid = str(uuid.uuid4())
            page = min(i + 1, pages)
            chunk_rows.append((cid, doc_id, i, page, chunk_text, cid))
            chroma_ids.append(cid)
            chroma_docs.append(chunk_text)
            chroma_metas.append({
                "doc_id": doc_id,
                "collection_id": row["collection_id"],
                "chunk_index": i,
                "page_number": page,
                "doc_type": row["doc_type"] or "unknown",
                "filename": row["filename"],
            })
            chroma_embeds.append(emb)

        db.executemany(
            "INSERT OR REPLACE INTO chunks VALUES (?,?,?,?,?,?)", chunk_rows
        )

        # Add to ChromaDB in batches
        batch = 50
        for b in range(0, len(chroma_ids), batch):
            sl = slice(b, b + batch)
            if all(e[0] != 0.0 for e in chroma_embeds[sl]):
                chroma_col.add(
                    ids=chroma_ids[sl],
                    documents=chroma_docs[sl],
                    metadatas=chroma_metas[sl],
                    embeddings=chroma_embeds[sl],
                )
            else:
                chroma_col.add(
                    ids=chroma_ids[sl],
                    documents=chroma_docs[sl],
                    metadatas=chroma_metas[sl],
                )

        # 5. Whoosh full-text index
        ix = _get_whoosh_index()
        writer = AsyncWriter(ix)
        for (cid, _, ci, page, chunk_text, _) in chunk_rows:
            writer.add_document(
                chunk_id=cid,
                doc_id=doc_id,
                collection_id=row["collection_id"],
                filename=row["filename"],
                title=row["title"] or row["filename"],
                body=chunk_text,
                manufacturer=row["manufacturer"] or "",
                part_number=row["part_number"] or "",
                doc_type=row["doc_type"] or "unknown",
                page_number=page,
            )
        writer.commit()

        # 6. NER
        entity_rows = _extract_entities(text[:50000])
        if entity_rows:
            db.executemany(
                "INSERT OR IGNORE INTO entities VALUES (?,?,?,?,?,?)",
                [(e["id"], doc_id, e["entity_type"], e["value"],
                  e["confidence"], e["context"]) for e in entity_rows],
            )

            # Update manufacturer / part_number on document if found
            mfrs = [e["value"] for e in entity_rows if e["entity_type"] == "MANUFACTURER"]
            parts = [e["value"] for e in entity_rows if e["entity_type"] == "PART_NUMBER"]
            if mfrs:
                db.execute("UPDATE documents SET manufacturer=? WHERE id=?",
                           (mfrs[0], doc_id))
            if parts:
                db.execute("UPDATE documents SET part_number=? WHERE id=?",
                           (parts[0], doc_id))

        db.execute(
            "UPDATE documents SET ingest_status='indexed', indexed_at=? WHERE id=?",
            (int(time.time()), doc_id),
        )
        db.commit()

    except Exception as exc:
        _set_status("error", str(exc)[:500])
    finally:
        db.close()
